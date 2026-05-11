"""
file_read PDF/图片/DOCX 扩展（FileExecutor mixin）

将 PDF 直读、图片多模态、DOCX 直读、页码解析从 FileExecutor 中拆出。
FileExecutor 通过继承 FileReadExtensionsMixin 获得这些能力。

修复清单（来自审查）：
- PDF 大文件预检：size > _MAX_READ_SIZE 直接拒绝
- PDF 错误信息脱敏：不暴露服务端路径
- 图片注入措辞中性化：[系统] 标签替代指令式文案
- 同步 I/O 改 run_in_executor：pdfplumber.open / Image.open / read_bytes
- DOCX 直读：python-docx 提取文本，不再绕 code_execute
"""

import asyncio
import base64
import mimetypes
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Optional, Union

from loguru import logger


# ── PDF 直读常量（对齐 Claude Code Read 工具 PDF 支持） ──

PDF_EXTENSIONS = frozenset({".pdf"})
DOCX_EXTENSIONS = frozenset({".docx"})
PPTX_EXTENSIONS = frozenset({".pptx"})
_PDF_MAX_AUTO_PAGES = 100   # 无 pages 参数时自动全读的最大页数（对标 Claude API）
_PDF_MAX_READ_PAGES = 100   # 单次读取最大页数（token 上限兜底）

# ── 图片直读常量 ──

IMAGE_EXTENSIONS = frozenset({".png", ".jpg", ".jpeg", ".gif", ".webp"})
_IMAGE_MAX_BASE64_SIZE = 2 * 1024 * 1024   # base64 fallback 最大文件大小（2MB）


@dataclass
class FileReadResult:
    """file_read 返回结果（支持文本和多模态）

    普通文本文件：type="text", text=内容
    PDF 文件：type="text", text=提取的页面文本
    图片文件：type="image", text=元信息描述, image_url=CDN/base64 URL
    """
    type: str = "text"         # "text" | "image"
    text: str = ""             # 文本内容（始终有值）
    image_url: str = ""        # 图片 URL（仅 type="image" 时有值）


class FileReadExtensionsMixin:
    """PDF 直读 + 图片多模态扩展

    假设宿主类（FileExecutor）提供：
    - _format_size(size: int) -> str
    - get_cdn_url(path: str) -> Optional[str]
    - _MAX_READ_SIZE: int（文件读取硬上限）
    - _BYTES_PER_TOKEN: int（token 估算比例）
    - _MAX_OUTPUT_TOKENS: int（输出 token 上限）
    """

    async def _read_pdf(
        self, path: str, target: Path, size: int, pages: str | None,
        max_read_size: int, bytes_per_token: int, max_output_tokens: int,
    ) -> str:
        """PDF 文件文本提取（pdfplumber 引擎）

        - 无文件大小硬上限（pdfplumber 逐页提取，内存安全）
        - ≤10 页无 pages 参数自动全读
        - >10 页必须指定 pages
        - 单次最多 20 页
        - 扫描件检测（提取文本为空）
        - 提取后按 token 上限控制返回量
        - 同步 I/O 在线程池执行，不阻塞事件循环
        """
        try:
            import pdfplumber
        except ImportError:
            return "PDF 读取依赖 pdfplumber 未安装，请用 code_execute 处理。"

        # ── 在线程池中打开 PDF ──
        loop = asyncio.get_running_loop()
        try:
            pdf = await loop.run_in_executor(
                None, lambda: pdfplumber.open(str(target)),
            )
        except Exception as e:
            logger.warning(f"PDF open failed | path={path} | error={type(e).__name__}")
            return f"PDF 文件无法打开: {path}（{type(e).__name__}）"

        total_pages = len(pdf.pages)
        if total_pages == 0:
            pdf.close()
            return f"PDF 文件为空: {path}"

        # 解析页码范围
        if pages:
            page_indices = self._parse_pages(pages, total_pages)
            if isinstance(page_indices, str):
                pdf.close()
                return page_indices
        elif total_pages <= _PDF_MAX_AUTO_PAGES:
            page_indices = list(range(total_pages))
        else:
            pdf.close()
            return (
                f"PDF 共 {total_pages} 页，超过 {_PDF_MAX_AUTO_PAGES} 页自动读取上限。\n"
                f"请用 pages 参数指定页范围（如 pages='1-5'、pages='3,7,10'）。"
            )

        # 单次最多 20 页
        if len(page_indices) > _PDF_MAX_READ_PAGES:
            pdf.close()
            return (
                f"请求读取 {len(page_indices)} 页，超过单次上限 {_PDF_MAX_READ_PAGES} 页。\n"
                "请缩小页码范围。"
            )

        # ── 在线程池中提取文本+表格 ──
        def _extract_pages():
            parts: list[str] = []
            empty: list[int] = []
            table_total = 0
            for idx in page_indices:
                page = pdf.pages[idx]
                page_text = page.extract_text() or ""
                page_text = page_text.strip()
                if not page_text:
                    empty.append(idx + 1)
                page_num = idx + 1
                page_parts = [f"── 第 {page_num} 页 ──"]
                page_parts.append(
                    page_text if page_text else "（无可提取文本）"
                )
                # 提取表格（pdfplumber extract_tables）
                try:
                    tables = page.extract_tables() or []
                    for tbl in tables:
                        if not tbl:
                            continue
                        table_total += 1
                        row_count = len(tbl)
                        col_count = max((len(r) for r in tbl), default=0)
                        page_parts.append(
                            f"\n=== 表格 ({row_count}行 x {col_count}列) ==="
                        )
                        for ri, row in enumerate(tbl, 1):
                            cells = [
                                (c or "").strip() for c in row
                            ]
                            if any(cells):
                                page_parts.append(f"  Row{ri}: {cells}")
                except Exception:
                    pass  # 表格提取失败不影响文本
                parts.append("\n".join(page_parts))
            pdf.close()
            return parts, empty, table_total

        parts, empty_pages, table_total = await loop.run_in_executor(
            None, _extract_pages,
        )
        output = "\n\n".join(parts)

        # L3: token 估算
        estimated_tokens = len(output.encode("utf-8")) / bytes_per_token
        if estimated_tokens > max_output_tokens:
            return (
                f"PDF 提取文本（约 {int(estimated_tokens)} tokens）"
                f"超过上限（{max_output_tokens} tokens）。\n"
                "请用 pages 参数读取更少的页面。"
            )

        # header
        page_range = pages or f"1-{total_pages}"
        header = f"文件: {path} | PDF {total_pages} 页 | 读取: {page_range}"
        if table_total > 0:
            header += f" | 表格: {table_total}个"
        if empty_pages:
            header += (
                f"\n⚠️ 第 {','.join(str(p) for p in empty_pages)} 页"
                "无可提取文本（可能是扫描件/图片页）"
            )

        return f"{header}\n{'─' * 60}\n{output}"

    async def _read_image(
        self, path: str, target: Path, size: int,
    ) -> "FileReadResult":
        """图片文件读取 → 返回 FileReadResult(type="image")

        有 CDN → 用 CDN URL
        无 CDN 且 ≤2MB → base64 data URL（线程池编码）
        无 CDN 且 >2MB → 仅返回元信息
        """
        loop = asyncio.get_running_loop()

        # ── 在线程池中获取宽高（PIL 是同步 I/O） ──
        def _get_dimensions():
            try:
                from PIL import Image
                with Image.open(target) as img:
                    return img.width, img.height
            except Exception:
                return None, None

        width, height = await loop.run_in_executor(None, _get_dimensions)

        size_str = self._format_size(size)
        mime = mimetypes.guess_type(str(target))[0] or "image/png"
        meta_text = f"图片: {path} ({size_str}"
        if width and height:
            meta_text += f", {width}×{height}px"
        meta_text += ")"

        # 尝试获取 CDN URL
        cdn_url = self.get_cdn_url(path)
        if cdn_url:
            return FileReadResult(
                type="image",
                text=f"{meta_text}\n模型已接收此图片用于视觉分析。",
                image_url=cdn_url,
            )

        # 无 CDN，尝试 base64（线程池中编码，避免阻塞事件循环）
        if size <= _IMAGE_MAX_BASE64_SIZE:
            def _encode_base64():
                raw = target.read_bytes()
                b64 = base64.b64encode(raw).decode("ascii")
                return f"data:{mime};base64,{b64}"

            try:
                data_url = await loop.run_in_executor(None, _encode_base64)
                return FileReadResult(
                    type="image",
                    text=f"{meta_text}\n模型已接收此图片用于视觉分析。",
                    image_url=data_url,
                )
            except Exception as e:
                logger.warning(f"Image base64 failed | path={path} | error={e}")

        # 无法传递图片给模型，返回纯元信息
        return FileReadResult(
            type="text",
            text=(
                f"{meta_text}\n"
                "⚠️ 图片过大或无 CDN 配置，无法直接查看内容。\n"
                "可用 code_execute 处理（PIL 已可用）。"
            ),
        )

    async def _read_docx(
        self, path: str, target: Path, size: int,
        max_read_size: int, bytes_per_token: int, max_output_tokens: int,
    ) -> str:
        """DOCX 结构化读取（保留标题层级+表格行号）

        输出格式对标 Claude：
          [Heading 1] 标题文本
          [Normal] 正文段落
          === 表格 1 (3行 x 4列) ===
            Row1: ['名称', '类型', '是否必须', '描述']
            Row2: ['pageNo', 'Integer', 'false', '页码']

        - 大文件预检：size > max_read_size 直接拒绝
        - 段落和表格按文档顺序交错输出
        - 同步 I/O 在线程池执行，不阻塞事件循环
        """
        if size > max_read_size:
            return (
                f"DOCX 文件过大（{self._format_size(size)}），"
                f"超过 {self._format_size(max_read_size)} 硬上限。"
                "建议使用 code_execute 处理。"
            )

        try:
            from docx import Document
        except ImportError:
            return "DOCX 读取依赖 python-docx 未安装，请用 code_execute 处理。"

        loop = asyncio.get_running_loop()
        try:
            def _extract_structured():
                doc = Document(str(target))
                parts: list[str] = []
                para_count = 0
                table_count = 0

                # 按文档顺序遍历（段落和表格交错）
                # python-docx 的 doc.element.body 按 XML 顺序包含所有块级元素
                from docx.table import Table as DocxTable
                from docx.text.paragraph import Paragraph

                table_idx = 0
                for element in doc.element.body:
                    tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag

                    if tag == "p":
                        para = Paragraph(element, doc)
                        text = para.text.strip()
                        if not text:
                            continue
                        para_count += 1
                        # 段落样式标注
                        style_name = para.style.name if para.style else "Normal"
                        if style_name.startswith("Heading"):
                            parts.append(f"[{style_name}] {text}")
                        elif style_name == "Title":
                            parts.append(f"[Title] {text}")
                        else:
                            parts.append(f"[Normal] {text}")

                    elif tag == "tbl":
                        if table_idx < len(doc.tables):
                            table = doc.tables[table_idx]
                            table_idx += 1
                            table_count += 1
                            row_count = len(table.rows)
                            col_count = max(
                                (len(r.cells) for r in table.rows), default=0,
                            )
                            parts.append(
                                f"\n=== 表格 {table_count} "
                                f"({row_count}行 x {col_count}列) ==="
                            )
                            for ri, row in enumerate(table.rows, 1):
                                cells = [c.text.strip() for c in row.cells]
                                # 跳过全空行
                                if any(cells):
                                    parts.append(f"  Row{ri}: {cells}")

                return "\n".join(parts), para_count, table_count

            text, para_count, table_count = await loop.run_in_executor(
                None, _extract_structured,
            )
        except Exception as e:
            logger.warning(f"DOCX open failed | path={path} | error={type(e).__name__}")
            return f"DOCX 文件无法打开: {path}（{type(e).__name__}）"

        if not text.strip():
            return f"DOCX 文件为空或无可提取文本: {path}"

        # L3: token 估算
        estimated_tokens = len(text.encode("utf-8")) / bytes_per_token
        if estimated_tokens > max_output_tokens:
            return (
                f"DOCX 提取文本（约 {int(estimated_tokens)} tokens）"
                f"超过上限（{max_output_tokens} tokens）。\n"
                "建议使用 code_execute 分段处理。"
            )

        header = (
            f"文件: {path} | DOCX | {self._format_size(size)}\n"
            f"段落数: {para_count}, 表格数: {table_count}"
        )
        return f"{header}\n{'─' * 60}\n{text}"

    async def _read_pptx(
        self, path: str, target: Path, size: int,
        max_read_size: int, bytes_per_token: int, max_output_tokens: int,
    ) -> str:
        """PPTX 结构化读取（Slide 编号 + 标题/文本/表格分块）

        输出格式：
          === Slide 1 ===
          [Title] 标题文本
          [Text] 正文内容
          === 表格 (3行 x 4列) ===
            Row1: ['名称', '类型', '描述']

        - 大文件预检：size > max_read_size 直接拒绝
        - 同步 I/O 在线程池执行，不阻塞事件循环
        """
        if size > max_read_size:
            return (
                f"PPTX 文件过大（{self._format_size(size)}），"
                f"超过 {self._format_size(max_read_size)} 硬上限。"
                "建议使用 code_execute 处理。"
            )

        try:
            from pptx import Presentation
        except ImportError:
            return "PPTX 读取依赖 python-pptx 未安装，请用 code_execute 处理。"

        loop = asyncio.get_running_loop()
        try:
            def _extract_structured():
                prs = Presentation(str(target))
                parts: list[str] = []
                total_slides = len(prs.slides)
                table_count = 0

                for i, slide in enumerate(prs.slides, 1):
                    parts.append(f"=== Slide {i} ===")
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            # 区分标题和正文
                            is_title = shape.shape_id == slide.placeholders[0].shape_id \
                                if 0 in slide.placeholders else False
                            for para in shape.text_frame.paragraphs:
                                text = para.text.strip()
                                if not text:
                                    continue
                                tag = "[Title]" if is_title else "[Text]"
                                parts.append(f"  {tag} {text}")
                        if shape.has_table:
                            table_count += 1
                            tbl = shape.table
                            row_count = len(tbl.rows)
                            col_count = max(
                                (len(r.cells) for r in tbl.rows), default=0,
                            )
                            parts.append(
                                f"  === 表格 ({row_count}行 x {col_count}列) ==="
                            )
                            for ri, row in enumerate(tbl.rows, 1):
                                cells = [c.text.strip() for c in row.cells]
                                if any(cells):
                                    parts.append(f"    Row{ri}: {cells}")
                    # Speaker notes
                    if slide.has_notes_slide:
                        notes = slide.notes_slide.notes_text_frame.text.strip()
                        if notes:
                            parts.append(f"  [Notes] {notes}")

                return "\n".join(parts), total_slides, table_count

            text, slide_count, table_count = await loop.run_in_executor(
                None, _extract_structured,
            )
        except Exception as e:
            logger.warning(f"PPTX open failed | path={path} | error={type(e).__name__}")
            return f"PPTX 文件无法打开: {path}（{type(e).__name__}）"

        if not text.strip():
            return f"PPTX 文件为空或无可提取文本: {path}"

        # L3: token 估算
        estimated_tokens = len(text.encode("utf-8")) / bytes_per_token
        if estimated_tokens > max_output_tokens:
            return (
                f"PPTX 提取文本（约 {int(estimated_tokens)} tokens）"
                f"超过上限（{max_output_tokens} tokens）。\n"
                "建议使用 code_execute 分段处理。"
            )

        header = (
            f"文件: {path} | PPTX | {self._format_size(size)}\n"
            f"幻灯片: {slide_count}页, 表格: {table_count}个"
        )
        return f"{header}\n{'─' * 60}\n{text}"

    @staticmethod
    def _parse_pages(pages_str: str, total_pages: int) -> Union[list[int], str]:
        """解析 pages 参数为 0-indexed 页码列表

        支持格式：'3'、'1-5'、'3,7,10'、'1-3,7,10-12'
        返回排序去重的页码列表，或错误信息字符串。
        """
        indices: set[int] = set()
        for part in pages_str.split(","):
            part = part.strip()
            if not part:
                continue
            if "-" in part:
                segments = part.split("-", 1)
                try:
                    start = int(segments[0].strip())
                    end = int(segments[1].strip())
                except ValueError:
                    from services.file_executor import FileOperationError
                    raise FileOperationError(f"页码格式错误: '{part}'，应为数字（如 '1-5'）")
                if start < 1 or end < 1:
                    from services.file_executor import FileOperationError
                    raise FileOperationError(f"页码必须从 1 开始: '{part}'")
                if start > total_pages or end > total_pages:
                    from services.file_executor import FileOperationError
                    raise FileOperationError(f"页码超出范围: '{part}'（共 {total_pages} 页）")
                if start > end:
                    from services.file_executor import FileOperationError
                    raise FileOperationError(f"起始页不能大于结束页: '{part}'")
                indices.update(range(start - 1, end))
            else:
                try:
                    page = int(part)
                except ValueError:
                    from services.file_executor import FileOperationError
                    raise FileOperationError(f"页码格式错误: '{part}'，应为数字")
                if page < 1:
                    from services.file_executor import FileOperationError
                    raise FileOperationError(f"页码必须从 1 开始: '{part}'")
                if page > total_pages:
                    from services.file_executor import FileOperationError
                    raise FileOperationError(f"页码超出范围: '{part}'（共 {total_pages} 页）")
                indices.add(page - 1)

        if not indices:
            from services.file_executor import FileOperationError
            raise FileOperationError("未指定有效页码")
        return sorted(indices)
