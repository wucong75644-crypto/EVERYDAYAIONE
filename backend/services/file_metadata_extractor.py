"""
文件元信息实时提取（对标 OpenAI Responses API spreadsheet augmentation）

上传工作区文件时，系统自动提取表格文件的元信息（行列数、类型、范围、分类），
注入 LLM context，让 AI 不用执行代码就知道文件结构。

设计文档：docs/document/TECH_文件元信息提取.md（待创建）
"""

import asyncio
import csv
import io
import os
import re
from collections import Counter
from datetime import datetime, date
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from loguru import logger


# ============================================================
# 常量
# ============================================================

_SPREADSHEET_EXTS = {".xlsx", ".xls", ".csv", ".tsv"}
_DOCUMENT_EXTS = {".docx", ".pptx", ".pdf"}
_TEXT_EXTS = {
    ".txt", ".md", ".log", ".json", ".jsonl", ".yaml", ".yml",
    ".xml", ".html", ".css", ".js", ".ts", ".py", ".sql",
    ".ini", ".cfg", ".conf", ".toml", ".rst", ".tex",
}
_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".svg"}
_DATA_EXTS = {".parquet"}
_SAMPLE_ROWS = 500  # 采样行数（对标 OpenAI 1000 行，ossfs 折中）
_MAX_PREVIEW_ROWS = 3  # 展示预览行数
_MAX_PREVIEW_COLS = 8  # 预览表格最多展示列数
_MAX_COLUMN_DISPLAY = 12  # 列名列表展示上限
_WIDE_TABLE_THRESHOLD = 50  # 列数超过此值尝试模式识别
_WIDE_TABLE_HEAD_COLS = 15  # 宽表退化模式：显示前 N 列名
_WIDE_TABLE_TAIL_COLS = 5  # 宽表退化模式：显示后 N 列名
_MAX_CATEGORY_DISPLAY = 5  # 分类值展示上限
_PER_FILE_TIMEOUT = 3.0  # 单文件超时（秒）
_TOTAL_TIMEOUT = 5.0  # 全部文件总超时（秒）
_MAX_METADATA_FILES = 5  # 最多提取文件数
_CSV_MAX_COUNT_SIZE = 50 * 1024 * 1024  # >50MB 跳过行数统计
_CATEGORY_THRESHOLD = 20  # unique ≤ 此值视为分类列


# ============================================================
# 类型推断
# ============================================================

def _is_numeric(value: Any) -> bool:
    """判断值是否为数值"""
    if isinstance(value, (int, float)):
        return True
    if isinstance(value, str):
        s = value.strip().replace(",", "").replace("，", "")
        if not s:
            return False
        try:
            float(s)
            return True
        except ValueError:
            return False
    return False


def _to_float(value: Any) -> Optional[float]:
    """尝试转为 float"""
    if isinstance(value, (int, float)):
        return float(value)
    if isinstance(value, str):
        s = value.strip().replace(",", "").replace("，", "")
        try:
            return float(s)
        except ValueError:
            return None
    return None


def _is_date(value: Any) -> bool:
    """判断值是否为日期"""
    if isinstance(value, (datetime, date)):
        return True
    if not isinstance(value, str):
        return False
    s = value.strip()
    if not s or len(s) < 6:
        return False
    # 常见日期格式
    for fmt in ("%Y-%m-%d", "%Y/%m/%d", "%Y.%m.%d", "%Y%m%d",
                "%Y-%m-%d %H:%M:%S", "%Y/%m/%d %H:%M:%S",
                "%m/%d/%Y", "%d/%m/%Y"):
        try:
            datetime.strptime(s[:19], fmt)
            return True
        except ValueError:
            continue
    return False


def _infer_column_meta(
    name: str,
    values: List[Any],
    total_rows: Optional[int],
) -> Dict[str, Any]:
    """从采样值推断单列元信息

    Returns:
        {
            "name": str,
            "dtype": "文本" | "数值" | "日期" | "布尔",
            "non_null": int,       # 估算非空数
            "sample": [str, ...],  # 前 2 个示例值
            "min": float,          # 数值列
            "max": float,          # 数值列
            "categories": [(val, count), ...],  # 分类列
        }
    """
    non_null_vals = [v for v in values if v is not None and str(v).strip() != ""]
    sample_size = len(values)
    non_null_count = len(non_null_vals)

    # 估算全量非空数
    if total_rows and sample_size > 0:
        estimated_non_null = round(total_rows * non_null_count / sample_size)
    else:
        estimated_non_null = non_null_count

    meta: Dict[str, Any] = {
        "name": name,
        "dtype": "文本",
        "non_null": estimated_non_null,
        "sample": [],
    }

    if not non_null_vals:
        return meta

    # 采集前 2 个不重复示例值
    seen = set()
    for v in non_null_vals:
        s = str(v).strip()
        if s and s not in seen:
            seen.add(s)
            meta["sample"].append(s)
            if len(meta["sample"]) >= 2:
                break

    # 类型推断
    numeric_count = sum(1 for v in non_null_vals if _is_numeric(v))
    date_count = sum(1 for v in non_null_vals if _is_date(v))
    bool_vals = {"true", "false", "是", "否", "0", "1", "yes", "no"}

    if numeric_count / len(non_null_vals) > 0.8:
        meta["dtype"] = "数值"
        # 数值范围
        floats = [f for f in (_to_float(v) for v in non_null_vals) if f is not None]
        if floats:
            meta["min"] = min(floats)
            meta["max"] = max(floats)
    elif date_count / len(non_null_vals) > 0.8:
        meta["dtype"] = "日期"
    elif all(str(v).lower().strip() in bool_vals for v in non_null_vals):
        meta["dtype"] = "布尔"

    # 分类检测（unique ≤ 阈值，且非数值/日期列）
    if meta["dtype"] == "文本":
        unique_vals = set(str(v).strip() for v in non_null_vals if str(v).strip())
        if 1 < len(unique_vals) <= _CATEGORY_THRESHOLD:
            counter = Counter(str(v).strip() for v in non_null_vals if str(v).strip())
            # 按出现次数降序，估算全量计数
            categories = []
            for val, cnt in counter.most_common(_MAX_CATEGORY_DISPLAY + 1):
                estimated = round(cnt * (total_rows or sample_size) / sample_size)
                categories.append((val, estimated))
            meta["categories"] = categories
            meta["_unique_count"] = len(unique_vals)

    return meta


# ============================================================
# 编码检测（复用 wecom/file_parser.py 模式）
# ============================================================

def _decode_bytes(data: bytes) -> str:
    """尝试 UTF-8 → GBK → Latin-1 解码"""
    for encoding in ("utf-8", "gbk", "latin-1"):
        try:
            return data.decode(encoding)
        except (UnicodeDecodeError, ValueError):
            continue
    return data.decode("utf-8", errors="replace")


# ============================================================
# 核心提取函数
# ============================================================

def extract_spreadsheet_metadata(abs_path: str) -> Optional[Dict[str, Any]]:
    """从表格文件提取元信息（同步，在线程池中执行）

    Args:
        abs_path: 文件绝对路径

    Returns:
        元信息 dict 或 None（提取失败）
    """
    if not os.path.exists(abs_path):
        return None

    ext = Path(abs_path).suffix.lower()
    try:
        if ext == ".xlsx":
            return _extract_xlsx(abs_path)
        elif ext == ".xls":
            # openpyxl 不支持 .xls (Excel 97-2003)，降级返回 None
            # 用户上传 .xls 时 LLM 仍可用 code_execute 读取
            logger.debug(f"Skipping .xls metadata extraction (openpyxl unsupported) | path={abs_path}")
            return None
        elif ext in (".csv", ".tsv"):
            # tsv 固定 tab，csv 自动检测分隔符
            delimiter = "\t" if ext == ".tsv" else _detect_csv_delimiter(abs_path)
            return _extract_csv(abs_path, delimiter=delimiter)
        return None
    except Exception as e:
        logger.warning(
            f"File metadata extraction failed | path={abs_path} | "
            f"error={type(e).__name__}: {e}"
        )
        return None


def _extract_xlsx(abs_path: str) -> Optional[Dict[str, Any]]:
    """提取 xlsx/xls 元信息"""
    from openpyxl import load_workbook

    wb = load_workbook(abs_path, read_only=True, data_only=True)
    try:
        # 提取所有 sheet 名称和行列数
        sheet_names = wb.sheetnames
        sheets_info = []
        for sn in sheet_names:
            s = wb[sn]
            sheets_info.append({
                "name": sn,
                "rows": (s.max_row or 0) - 1 if s.max_row and s.max_row > 0 else 0,
                "cols": s.max_column or 0,
            })

        # 选第一个 sheet 做详细提取
        ws = wb[sheet_names[0]]
        total_rows = ws.max_row
        total_cols = ws.max_column

        # 减去表头行
        data_rows = (total_rows - 1) if total_rows and total_rows > 0 else 0

        # 读取采样行（表头 + 数据行）
        max_read = min(_SAMPLE_ROWS + 1, (total_rows or _SAMPLE_ROWS + 1))
        all_rows = []
        for row in ws.iter_rows(min_row=1, max_row=max_read, values_only=True):
            all_rows.append(list(row))
            if len(all_rows) >= max_read:
                break

        if not all_rows:
            return None

        # 表头检测：messytables 列数众数法 + csv.Sniffer 类型验证
        # 典型场景：ERP 利润表前几行是标题/分类行，真正表头在第 2-3 行
        from services.agent.data_query_cache import detect_header_row
        header_row_idx = detect_header_row(all_rows)

        header = [str(c).strip() if c is not None else f"列{i+1}"
                  for i, c in enumerate(all_rows[header_row_idx])]
        data = all_rows[header_row_idx + 1:]

        # 修正数据行数（减去实际表头行号）
        data_rows = max(0, (total_rows or 0) - header_row_idx - 1)

        # 按列聚合
        columns_meta = []
        for col_idx, col_name in enumerate(header):
            col_values = [row[col_idx] if col_idx < len(row) else None for row in data]
            columns_meta.append(_infer_column_meta(col_name, col_values, data_rows))

        # 预览行
        preview_rows = []
        for row in data[:_MAX_PREVIEW_ROWS]:
            preview_rows.append([
                _format_cell(row[i] if i < len(row) else None)
                for i in range(min(len(header), _MAX_PREVIEW_COLS))
            ])

        result = {
            "row_count": data_rows,
            "col_count": total_cols or len(header),
            "columns": columns_meta,
            "preview_rows": preview_rows,
        }
        # 非首行表头 → 记录 header 参数，让读取指引自动带上
        if header_row_idx > 0:
            result["header_row"] = header_row_idx
        # 多 sheet 信息（仅多 sheet 时注入，单 sheet 不展示）
        if len(sheets_info) > 1:
            result["sheets"] = sheets_info
        return result
    finally:
        wb.close()


# ============================================================
# 文档类提取（docx / pptx / pdf）
# ============================================================

_MAX_DOC_PREVIEW_CHARS = 500  # 文档预览文本最大字符数
_MAX_DOC_PREVIEW_PARAGRAPHS = 10  # 预览最多段落数


def extract_document_metadata(abs_path: str) -> Optional[Dict[str, Any]]:
    """从文档文件提取元信息（同步，在线程池中执行）"""
    if not os.path.exists(abs_path):
        return None

    ext = Path(abs_path).suffix.lower()
    try:
        if ext == ".docx":
            return _extract_docx(abs_path)
        elif ext == ".pptx":
            return _extract_pptx(abs_path)
        elif ext == ".pdf":
            return _extract_pdf(abs_path)
        return None
    except Exception as e:
        logger.warning(
            f"Document metadata extraction failed | path={abs_path} | "
            f"error={type(e).__name__}: {e}"
        )
        return None


def _extract_docx(abs_path: str) -> Optional[Dict[str, Any]]:
    """提取 docx 元信息：段落数、表格数、前几段预览"""
    from docx import Document

    doc = Document(abs_path)
    paragraphs = doc.paragraphs
    tables = doc.tables

    # 提取非空段落文本
    non_empty = [p.text.strip() for p in paragraphs if p.text.strip()]
    total_paragraphs = len(non_empty)

    # 预览前 N 段
    preview_parts = []
    char_count = 0
    for text in non_empty[:_MAX_DOC_PREVIEW_PARAGRAPHS]:
        if char_count + len(text) > _MAX_DOC_PREVIEW_CHARS:
            remaining = _MAX_DOC_PREVIEW_CHARS - char_count
            if remaining > 20:
                preview_parts.append(text[:remaining] + "...")
            break
        preview_parts.append(text)
        char_count += len(text)

    # 总字数估算
    total_chars = sum(len(t) for t in non_empty)

    return {
        "type": "docx",
        "paragraphs": total_paragraphs,
        "tables": len(tables),
        "chars": total_chars,
        "preview": preview_parts,
    }


def _extract_pptx(abs_path: str) -> Optional[Dict[str, Any]]:
    """提取 pptx 元信息：幻灯片数、每页标题、总字数"""
    from pptx import Presentation

    prs = Presentation(abs_path)
    slides = prs.slides
    total_slides = len(slides)

    slide_info = []
    total_chars = 0
    for i, slide in enumerate(slides):
        title = ""
        slide_text_parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    slide_text_parts.append(text)
                    total_chars += len(text)
            if hasattr(shape, "text") and shape == slide.shapes.title:
                title = shape.text.strip()

        slide_info.append({
            "index": i + 1,
            "title": title or f"第{i+1}页",
            "text_len": sum(len(t) for t in slide_text_parts),
        })

    return {
        "type": "pptx",
        "slides": total_slides,
        "slide_titles": [s["title"] for s in slide_info[:20]],  # 最多 20 页标题
        "chars": total_chars,
    }


def _extract_pdf(abs_path: str) -> Optional[Dict[str, Any]]:
    """提取 pdf 元信息：页数、前 1 页文本预览"""
    try:
        from PyPDF2 import PdfReader
    except ImportError:
        logger.debug("PyPDF2 not installed, skipping PDF metadata")
        return None

    reader = PdfReader(abs_path)
    total_pages = len(reader.pages)

    # 提取前 1 页文本预览
    preview = ""
    if total_pages > 0:
        first_page_text = reader.pages[0].extract_text() or ""
        preview = first_page_text.strip()[:_MAX_DOC_PREVIEW_CHARS]
        if len(first_page_text.strip()) > _MAX_DOC_PREVIEW_CHARS:
            preview += "..."

    # 估算总字数（采样前 3 页）
    total_chars = 0
    sample_pages = min(3, total_pages)
    for i in range(sample_pages):
        page_text = reader.pages[i].extract_text() or ""
        total_chars += len(page_text)
    if sample_pages > 0 and total_pages > sample_pages:
        total_chars = round(total_chars / sample_pages * total_pages)

    result = {
        "type": "pdf",
        "pages": total_pages,
        "chars": total_chars,
        "preview": preview,
    }
    # 扫描件检测：有页但无文本 → 标注为扫描 PDF
    if total_pages > 0 and total_chars < 10:
        result["is_scanned"] = True
    return result


def _detect_csv_delimiter(abs_path: str) -> str:
    """自动检测 CSV 分隔符（csv.Sniffer + 兜底逗号）"""
    try:
        with open(abs_path, "rb") as f:
            head = f.read(8192)  # 前 8KB 足够判断
        text = _decode_bytes(head)
        dialect = csv.Sniffer().sniff(text, delimiters=",;\t|")
        return dialect.delimiter
    except Exception:
        return ","


def _extract_csv(abs_path: str, delimiter: str = ",") -> Optional[Dict[str, Any]]:
    """提取 csv/tsv 元信息

    单次读取策略：流式扫描文件，前 500 行进入采样，之后只计数不存储。
    避免 ossfs 上的双次 IO 延迟。
    """
    file_size = os.path.getsize(abs_path)

    # 读取整个文件的原始字节（编码检测需要）
    # 对于大文件（>50MB），只读头部用于采样，跳过行数统计
    skip_count = file_size > _CSV_MAX_COUNT_SIZE
    if skip_count:
        with open(abs_path, "rb") as f:
            raw = f.read(2 * 1024 * 1024)  # 2MB 足够覆盖 500 行宽表
    else:
        with open(abs_path, "rb") as f:
            raw = f.read()

    text = _decode_bytes(raw)
    reader = csv.reader(io.StringIO(text), delimiter=delimiter)

    # 单次扫描：前 _SAMPLE_ROWS+1 行进入采样，之后只计数
    sample_rows: List[List[str]] = []
    total_line_count = 0
    for row in reader:
        total_line_count += 1
        if total_line_count <= _SAMPLE_ROWS + 1:
            sample_rows.append(row)

    if not sample_rows:
        return None

    header = [c.strip() if c.strip() else f"列{i+1}" for i, c in enumerate(sample_rows[0])]
    data = sample_rows[1:]

    # 行数：total_line_count 包含表头，减 1 得数据行数
    # 大文件跳过计数时，如果采样不足 _SAMPLE_ROWS 说明 2MB 内全部读完
    total_data_rows: Optional[int] = None
    if not skip_count:
        total_data_rows = max(0, total_line_count - 1)
    elif len(data) < _SAMPLE_ROWS:
        total_data_rows = len(data)

    # 按列聚合
    columns_meta = []
    for col_idx, col_name in enumerate(header):
        col_values = [row[col_idx] if col_idx < len(row) else None for row in data]
        columns_meta.append(_infer_column_meta(col_name, col_values, total_data_rows))

    # 预览行
    preview_rows = []
    for row in data[:_MAX_PREVIEW_ROWS]:
        preview_rows.append([
            _format_cell(row[i] if i < len(row) else None)
            for i in range(min(len(header), _MAX_PREVIEW_COLS))
        ])

    result = {
        "row_count": total_data_rows,
        "col_count": len(header),
        "columns": columns_meta,
        "preview_rows": preview_rows,
    }
    # 非逗号分隔 → 记录分隔符，让读取指引自动带 sep 参数
    if delimiter != ",":
        _SEP_NAMES = {"\t": "\\t", ";": ";", "|": "|"}
        result["delimiter"] = _SEP_NAMES.get(delimiter, delimiter)
    return result


def _format_cell(value: Any) -> str:
    """格式化单元格值用于预览展示"""
    if value is None:
        return ""
    s = str(value).strip()
    # 截断过长的值
    if len(s) > 30:
        return s[:27] + "..."
    return s


# ============================================================
# 异步批量包装
# ============================================================

async def extract_metadata_for_files(
    workspace_files: List[Dict[str, Any]],
    workspace_dir: str,
) -> Dict[str, Optional[Dict[str, Any]]]:
    """批量提取文件元信息（并行 + 超时保护）

    Args:
        workspace_files: workspace 文件列表 [{workspace_path, name, size, ...}]
        workspace_dir: 用户 workspace 绝对路径

    Returns:
        {workspace_path: metadata_dict_or_None}
    """
    result: Dict[str, Optional[Dict[str, Any]]] = {}

    # 所有文件都可提取元数据（统一入口 extract_file_metadata 按扩展名分派）
    extractable = list(workspace_files)[:_MAX_METADATA_FILES]

    if not extractable:
        return result

    loop = asyncio.get_running_loop()

    async def _extract_one(ws_file: Dict[str, Any]) -> Tuple[str, Optional[Dict[str, Any]]]:
        wp = ws_file["workspace_path"]
        abs_path = os.path.join(workspace_dir, wp)
        try:
            meta = await asyncio.wait_for(
                loop.run_in_executor(None, extract_file_metadata, abs_path),
                timeout=_PER_FILE_TIMEOUT,
            )
            return wp, meta
        except Exception as e:
            logger.debug(f"Metadata extraction skipped | path={wp} | error={e}")
            return wp, None

    try:
        tasks = [_extract_one(f) for f in extractable]
        results = await asyncio.wait_for(
            asyncio.gather(*tasks, return_exceptions=True),
            timeout=_TOTAL_TIMEOUT,
        )
        for r in results:
            if isinstance(r, tuple):
                result[r[0]] = r[1]
            # BaseException 忽略
    except asyncio.TimeoutError:
        logger.warning(
            f"File metadata extraction global timeout | "
            f"files={[f['workspace_path'] for f in spreadsheets]}"
        )

    return result


# ============================================================
# 格式化输出
# ============================================================

def format_workspace_files_prompt(
    workspace_files: List[Dict[str, Any]],
    metadata_map: Dict[str, Optional[Dict[str, Any]]],
) -> str:
    """将工作区文件信息格式化为 LLM system prompt

    只展示靠谱的静态信息：文件名、大小、类型、时间、路径。
    不提取列名、不检测表头、不做模式识别——让 AI 用 code_execute 自己看。
    """
    if not workspace_files:
        return ""

    parts = ["用户工作区文件：\n"]

    for f in workspace_files:
        wp = f.get("workspace_path", "")
        ext = Path(wp).suffix.lower() if wp else ""
        size_str = _fmt_size(f.get("size"))
        mtime = f.get("modified", "")

        # 文件类型图标
        if ext in _SPREADSHEET_EXTS:
            icon = "📊"
        elif ext in _DOCUMENT_EXTS:
            icon = "📝"
        elif ext in {".png", ".jpg", ".jpeg", ".gif", ".webp", ".svg"}:
            icon = "🖼️"
        else:
            icon = "📄"

        parts.append(f"{icon} {wp} | {size_str} | {ext.lstrip('.')} | {mtime}")
        parts.append(f"   路径: {wp}")

    return "\n".join(parts)



# ============================================================
# 宽表模式识别
# ============================================================

# 日期模式：2024-01, 2024/01, 2024年1月, 202401, 1月, 01 等
_DATE_PATTERNS = [
    re.compile(r"\d{4}[-/]\d{1,2}"),        # 2024-01, 2024/1
    re.compile(r"\d{4}年\d{1,2}月"),          # 2024年1月
    re.compile(r"\d{6}"),                     # 202401
    re.compile(r"\d{1,2}月"),                 # 1月, 12月
]


def _detect_wide_table_pattern(
    columns: List[Dict[str, Any]],
) -> Optional[Dict[str, Any]]:
    """检测宽表列名模式（列数 > _WIDE_TABLE_THRESHOLD 时调用）

    返回 None 表示无法识别模式，调用方应退化为列名列表。
    返回 dict 包含：
      - description: 模式描述文本
      - index_cols: 行索引列信息
      - value_col_count: 值列数量
    """
    if len(columns) <= _WIDE_TABLE_THRESHOLD:
        return None

    col_names = [c["name"] for c in columns]

    # 1. 找分隔符：从常用分隔符中选命中率最高的
    best_sep, best_count = None, 0
    for sep in ["_", "-", "/", "｜", "|"]:
        count = sum(1 for n in col_names if sep in n)
        if count > best_count:
            best_sep, best_count = sep, count

    # 分隔符命中率 < 60% → 无法识别模式
    if best_sep is None or best_count < len(col_names) * 0.6:
        return None

    # 2. 按分隔符拆分列名，分离索引列和值列
    prefixes: List[str] = []
    suffixes: List[str] = []
    index_cols: List[Dict[str, Any]] = []

    for col in columns:
        name = col["name"]
        if best_sep not in name:
            index_cols.append(col)
            continue
        parts = name.split(best_sep, 1)
        prefixes.append(parts[0].strip())
        suffixes.append(parts[1].strip() if len(parts) > 1 else "")

    if not prefixes:
        return None

    # 3. 检测后缀是否为日期模式
    is_date_suffix = False
    for pattern in _DATE_PATTERNS:
        date_hits = sum(1 for s in suffixes if pattern.search(s))
        if date_hits > len(suffixes) * 0.5:
            is_date_suffix = True
            break

    # 4. 聚合统计
    unique_prefixes = sorted(set(prefixes))
    unique_suffixes = sorted(set(suffixes))

    # 5. 构建模式描述
    # 索引列描述
    idx_desc = ""
    if index_cols:
        idx_names = [c["name"] for c in index_cols[:5]]
        idx_desc = f"行索引列: {', '.join(idx_names)}"
        if len(index_cols) > 5:
            idx_desc += f" 等共 {len(index_cols)} 列"

    # 值列模式
    suffix_label = "日期" if is_date_suffix else "后缀"
    prefix_sample = unique_prefixes[:3]
    suffix_sample = unique_suffixes[:3]

    # 示例列名
    sample_cols = [col_names[i] for i in range(len(col_names))
                   if best_sep in col_names[i]][:3]

    description_lines = [
        f"结构: 宽表（{len(unique_prefixes)} 个前缀 × {len(unique_suffixes)} 个{suffix_label}）",
    ]
    if idx_desc:
        description_lines.append(idx_desc)
    description_lines.append(
        f"值列模式: {{前缀}}{best_sep}{{{suffix_label}}}，共 {len(prefixes)} 列"
    )
    description_lines.append(
        f"  前缀({len(unique_prefixes)}): {', '.join(prefix_sample)}"
        + (f", ..." if len(unique_prefixes) > 3 else "")
    )
    description_lines.append(
        f"  {suffix_label}({len(unique_suffixes)}): {', '.join(suffix_sample)}"
        + (f", ..." if len(unique_suffixes) > 3 else "")
    )
    description_lines.append(
        f"  示例列: {', '.join(sample_cols)}"
    )

    # 值列数据类型（用 id 集合避免 O(n²) dict 比较）
    _index_ids = set(id(c) for c in index_cols)
    value_cols = [c for c in columns if id(c) not in _index_ids]
    if value_cols:
        dtypes = Counter(c["dtype"] for c in value_cols)
        main_dtype = dtypes.most_common(1)[0][0]
        description_lines.append(f"  值类型: {main_dtype}")
        # 如果有范围信息，取所有值列的全局 min/max
        mins = [c["min"] for c in value_cols if c.get("min") is not None]
        maxs = [c["max"] for c in value_cols if c.get("max") is not None]
        if mins and maxs:
            global_min = min(mins)
            global_max = max(maxs)
            min_s = f"{global_min:,.0f}" if isinstance(global_min, (int, float)) and float(global_min).is_integer() else f"{global_min:,.2f}"
            max_s = f"{global_max:,.0f}" if isinstance(global_max, (int, float)) and float(global_max).is_integer() else f"{global_max:,.2f}"
            description_lines.append(f"  值范围: {min_s} ~ {max_s}")

    return {
        "description": "\n".join(description_lines),
        "index_cols": index_cols,
        "value_col_count": len(prefixes),
    }


def _format_wide_table_fallback(columns: List[Dict[str, Any]]) -> str:
    """宽表模式识别失败时的退化方案：展示前 N + 后 N 列名"""
    head = columns[:_WIDE_TABLE_HEAD_COLS]
    tail = columns[-_WIDE_TABLE_TAIL_COLS:] if len(columns) > _WIDE_TABLE_HEAD_COLS + _WIDE_TABLE_TAIL_COLS else []
    parts = [c["name"] for c in head]
    if tail:
        parts.append(f"... (省略 {len(columns) - _WIDE_TABLE_HEAD_COLS - _WIDE_TABLE_TAIL_COLS} 列)")
        parts.extend(c["name"] for c in tail)
    return "列名: " + ", ".join(parts)


def _format_standard(
    wp: str, size_str: str, meta: Dict[str, Any]
) -> str:
    """标准档格式化（详细表格，~200-300 tokens）"""
    row_count = meta.get("row_count")
    col_count = meta.get("col_count", 0)
    columns = meta.get("columns", [])
    preview_rows = meta.get("preview_rows", [])

    row_label = f"{row_count:,}" if row_count is not None else "未知"

    lines = [
        f"📊 {wp} | {row_label} 行 × {col_count} 列",
        f"路径: {wp}",
        "",
    ]

    # 宽表检测：列数 > 阈值时尝试模式识别
    if col_count > _WIDE_TABLE_THRESHOLD and columns:
        pattern = _detect_wide_table_pattern(columns)
        if pattern:
            lines.append(pattern["description"])
        else:
            lines.append(_format_wide_table_fallback(columns))
    else:
        # 常规表：列元信息表格（最多 _MAX_PREVIEW_COLS 列）
        display_cols = columns[:_MAX_PREVIEW_COLS]
        if display_cols:
            lines.append("| 列名 | 类型 | 非空 | 示例值 |")
            lines.append("|------|------|------|--------|")

            for col in display_cols:
                name = col["name"]
                dtype = col["dtype"]
                non_null = col.get("non_null", "")
                non_null_str = f"{non_null:,}" if isinstance(non_null, int) else str(non_null)

                # 构建示例值
                sample_parts = []
                if col.get("categories"):
                    cats = col["categories"]
                    for val, cnt in cats[:_MAX_CATEGORY_DISPLAY]:
                        sample_parts.append(f'"{val}"({cnt:,})')
                    unique = col.get("_unique_count", len(cats))
                    if unique > _MAX_CATEGORY_DISPLAY:
                        sample_parts.append(f"+{unique - _MAX_CATEGORY_DISPLAY}类")
                elif col.get("sample"):
                    for s in col["sample"]:
                        sample_parts.append(f'"{s}"' if dtype != "数值" else s)
                    if col.get("min") is not None and col.get("max") is not None:
                        mn, mx = col["min"], col["max"]
                        mn_s = f"{mn:,.0f}" if float(mn).is_integer() else f"{mn:,.2f}"
                        mx_s = f"{mx:,.0f}" if float(mx).is_integer() else f"{mx:,.2f}"
                        sample_parts.append(f"[范围: {mn_s}~{mx_s}]")

                sample_str = ", ".join(sample_parts) if sample_parts else ""
                lines.append(f"| {name} | {dtype} | {non_null_str} | {sample_str} |")

            if len(columns) > _MAX_PREVIEW_COLS:
                lines.append(f"| ... | | | (共{col_count}列，显示前{_MAX_PREVIEW_COLS}列) |")

    # 多 sheet 提示
    sheets = meta.get("sheets")
    if sheets:
        sheet_parts = [f"{s['name']}({s['rows']}行)" for s in sheets]
        lines.append(f"\n📑 多 Sheet: {', '.join(sheet_parts)}")
        lines.append(f"  以上为第一个 Sheet「{sheets[0]['name']}」的信息。"
                     f"读其他 Sheet: data_query(file='{wp}', sheet='Sheet名')")

    # 读取指引（自动带正确参数）
    read_cmd = _build_read_command(wp, meta)
    if read_cmd:
        lines.append(f"\n用 data_query 分析: {read_cmd}")

    return "\n".join(lines)


def _format_compact(
    wp: str, size_str: str, meta: Dict[str, Any]
) -> str:
    """紧凑档格式化（单行摘要，~80 tokens）"""
    row_count = meta.get("row_count")
    col_count = meta.get("col_count", 0)
    columns = meta.get("columns", [])

    row_label = f"{row_count:,}" if row_count is not None else "?"

    # 列摘要
    col_parts = []
    for col in columns[:_MAX_COLUMN_DISPLAY]:
        name = col["name"]
        dtype = col["dtype"]
        extra = ""
        if col.get("min") is not None and col.get("max") is not None:
            mn, mx = col["min"], col["max"]
            mn_s = f"{mn:,.0f}" if float(mn).is_integer() else f"{mn:.1f}"
            mx_s = f"{mx:,.0f}" if float(mx).is_integer() else f"{mx:.1f}"
            extra = f",{mn_s}~{mx_s}"
        elif col.get("_unique_count"):
            extra = f",{col['_unique_count']}类"

        # 标注空值
        non_null = col.get("non_null", 0)
        row_total = row_count or 0
        null_count = row_total - non_null if row_total and non_null else 0
        if null_count > 0:
            extra += f",{null_count}空"

        col_parts.append(f"{name}({dtype}{extra})")

    cols_str = ", ".join(col_parts)
    if len(columns) > _MAX_COLUMN_DISPLAY:
        cols_str += ", ..."

    # 多 sheet 紧凑提示
    sheets = meta.get("sheets")
    sheet_hint = ""
    if sheets:
        sheet_names = [s["name"] for s in sheets]
        sheet_hint = f"\n  📑 {len(sheets)} Sheets: {', '.join(sheet_names)}"

    read_cmd = _build_read_command(wp, meta)

    return (
        f"📊 {wp} | {row_label}×{col_count} | 路径: {wp}\n"
        f"  列: {cols_str}{sheet_hint}\n"
        f"  用 {read_cmd} 读取"
    )


def _build_read_command(wp: str, meta: Dict[str, Any]) -> str:
    """根据文件类型和元信息生成 data_query 读取命令

    核心设计：元信息提取发现的特殊情况（Sheet 名）
    自动反映到读取命令中，AI 复制即可用，不需要自己猜参数。
    表头行由 data_query 内部自动检测，无需暴露给 LLM。
    """
    params = [f"file=\"{wp}\""]

    # 多 sheet 时提示 sheet 参数
    sheets = meta.get("sheets")
    if sheets and len(sheets) > 1:
        params.append(f"sheet=\"{sheets[0]['name']}\"")

    return f"data_query({', '.join(params)})"


def _format_document_entry(
    wp: str, size_str: str, ext: str, meta: Optional[Dict[str, Any]],
) -> str:
    """格式化文档文件条目（docx/pptx/pdf）"""
    _DOC_READ_HINTS = {
        ".docx": f"file_read(path=\"{wp}\")",
        ".pptx": f"from pptx import Presentation; prs = Presentation('{wp}')",
        ".pdf": f"file_read(path=\"{wp}\")",
    }

    read_hint = _DOC_READ_HINTS.get(ext, f"open('{wp}')")

    if meta is None:
        action = "读取" if ext in {".docx", ".pdf"} else "用 code_execute"
        return f"📄 {wp} ({size_str})\n  {action}: {read_hint}"

    doc_type = meta.get("type", "")

    if doc_type == "docx":
        paragraphs = meta.get("paragraphs", 0)
        tables = meta.get("tables", 0)
        chars = meta.get("chars", 0)
        preview = meta.get("preview", [])

        lines = [f"📝 {wp} ({size_str}) | {paragraphs}段 {tables}表 ~{chars:,}字"]
        if preview:
            preview_text = " / ".join(preview[:3])
            if len(preview_text) > 100:
                preview_text = preview_text[:97] + "..."
            lines.append(f"  预览: {preview_text}")
        lines.append(f"  读取: {read_hint}")
        return "\n".join(lines)

    elif doc_type == "pptx":
        slides = meta.get("slides", 0)
        chars = meta.get("chars", 0)
        titles = meta.get("slide_titles", [])

        lines = [f"📽 {wp} ({size_str}) | {slides}页 ~{chars:,}字"]
        if titles:
            title_preview = ", ".join(titles[:5])
            if len(titles) > 5:
                title_preview += f" (+{len(titles)-5}页)"
            lines.append(f"  页标题: {title_preview}")
        lines.append(f"  用 code_execute: {read_hint}")
        return "\n".join(lines)

    elif doc_type == "pdf":
        pages = meta.get("pages", 0)
        chars = meta.get("chars", 0)
        preview = meta.get("preview", "")
        is_scanned = meta.get("is_scanned", False)

        lines = [f"📕 {wp} ({size_str}) | {pages}页 ~{chars:,}字"]
        if is_scanned:
            lines.append("  ⚠️ 扫描件 PDF（无可提取文本），需 OCR 处理")
        elif preview:
            short = preview[:100] + "..." if len(preview) > 100 else preview
            lines.append(f"  首页预览: {short}")
        lines.append(f"  用 code_execute: {read_hint}")
        return "\n".join(lines)

    return f"📄 {wp} ({size_str})\n  用 code_execute: {read_hint}"


# ============================================================
# 文本文件 / 图片 / Parquet 提取
# ============================================================

_TEXT_PREVIEW_LINES = 5
_TEXT_MAX_PREVIEW_CHARS = 500


def _extract_text_file(abs_path: str) -> Optional[Dict[str, Any]]:
    """提取文本文件元信息：行数、字数、前几行预览"""
    file_size = os.path.getsize(abs_path)
    if file_size > 10 * 1024 * 1024:  # >10MB 跳过
        return {"type": "text", "lines": None, "chars": None, "preview": "(文件过大)"}

    try:
        with open(abs_path, "rb") as f:
            raw = f.read()
        text = _decode_bytes(raw)
    except Exception:
        return None

    lines = text.splitlines()
    total_lines = len(lines)
    total_chars = len(text)

    # 预览前 N 行
    preview_lines = []
    char_count = 0
    for line in lines[:_TEXT_PREVIEW_LINES]:
        if char_count + len(line) > _TEXT_MAX_PREVIEW_CHARS:
            remaining = _TEXT_MAX_PREVIEW_CHARS - char_count
            if remaining > 20:
                preview_lines.append(line[:remaining] + "...")
            break
        preview_lines.append(line)
        char_count += len(line)

    return {
        "type": "text",
        "lines": total_lines,
        "chars": total_chars,
        "preview": preview_lines,
    }


def _extract_image(abs_path: str) -> Optional[Dict[str, Any]]:
    """提取图片元信息：宽高像素"""
    try:
        from PIL import Image
        with Image.open(abs_path) as img:
            return {
                "type": "image",
                "width": img.width,
                "height": img.height,
                "mode": img.mode,
            }
    except Exception:
        return None


def _extract_parquet(abs_path: str) -> Optional[Dict[str, Any]]:
    """提取 parquet 元信息：行列数、列名"""
    try:
        import pyarrow.parquet as pq
        pf = pq.ParquetFile(abs_path)
        schema = pf.schema_arrow
        num_rows = pf.metadata.num_rows
        col_names = [field.name for field in schema]

        return {
            "type": "parquet",
            "row_count": num_rows,
            "col_count": len(col_names),
            "columns": col_names[:20],  # 最多展示 20 列名
        }
    except Exception:
        return None


# ============================================================
# 统一提取入口（所有文件类型）
# ============================================================

def extract_file_metadata(abs_path: str) -> Optional[Dict[str, Any]]:
    """统一文件元信息提取入口

    根据文件扩展名选择对应的提取函数。
    file_list/file_search 统一调用此函数。
    """
    if not os.path.exists(abs_path):
        return None

    ext = Path(abs_path).suffix.lower()

    if ext in _SPREADSHEET_EXTS:
        return extract_spreadsheet_metadata(abs_path)
    elif ext in _DOCUMENT_EXTS:
        return extract_document_metadata(abs_path)
    elif ext in _TEXT_EXTS:
        return _extract_text_file(abs_path)
    elif ext in _IMAGE_EXTS:
        return _extract_image(abs_path)
    elif ext in _DATA_EXTS:
        return _extract_parquet(abs_path)
    return None


def format_file_metadata_line(
    name: str, abs_path: str, size: int, meta: Optional[Dict[str, Any]],
) -> str:
    """将单个文件的元数据格式化为一行展示文本

    供 file_list/file_search 的返回结果使用。
    """
    size_str = _fmt_size(size)
    ext = Path(name).suffix.lower()

    if meta is None:
        return f"  {name}\t{size_str}"

    file_type = meta.get("type", "")

    # 表格类
    if file_type in ("", ) or ext in _SPREADSHEET_EXTS:
        row_count = meta.get("row_count")
        col_count = meta.get("col_count", 0)
        row_label = f"{row_count:,}" if row_count is not None else "?"
        sheets = meta.get("sheets")
        sheet_hint = f" | {len(sheets)} Sheets" if sheets else ""
        header_row = meta.get("header_row")
        read_cmd = _build_read_command(name, meta)
        header_hint = f" | ⚠️ header={header_row}" if header_row else ""
        return (
            f"  📊 {name}\t{size_str} | {row_label}行×{col_count}列"
            f"{sheet_hint}{header_hint}\n"
            f"     读取: {read_cmd}"
        )

    # 文档类
    if file_type == "docx":
        p = meta.get("paragraphs", 0)
        t = meta.get("tables", 0)
        c = meta.get("chars", 0)
        return (
            f"  📝 {name}\t{size_str} | {p}段 {t}表 ~{c:,}字\n"
            f"     读取: file_read(path=\"{name}\")"
        )
    if file_type == "pptx":
        s = meta.get("slides", 0)
        c = meta.get("chars", 0)
        titles = meta.get("slide_titles", [])
        title_hint = f" | {', '.join(titles[:3])}" if titles else ""
        return (
            f"  📽 {name}\t{size_str} | {s}页 ~{c:,}字{title_hint}\n"
            f"     读取: file_read(path=\"{name}\")"
        )
    if file_type == "pdf":
        p = meta.get("pages", 0)
        c = meta.get("chars", 0)
        scanned = " | ⚠️扫描件" if meta.get("is_scanned") else ""
        return (
            f"  📕 {name}\t{size_str} | {p}页 ~{c:,}字{scanned}\n"
            f"     读取: file_read(path=\"{name}\")"
        )

    # 文本类
    if file_type == "text":
        lines = meta.get("lines")
        chars = meta.get("chars")
        lines_label = f"{lines:,}行" if lines else "?"
        chars_label = f"~{chars:,}字" if chars else ""
        preview = meta.get("preview", [])
        first_line = preview[0][:60] + "..." if preview and len(preview[0]) > 60 else (preview[0] if preview else "")
        preview_hint = f' | "{first_line}"' if first_line else ""
        return (
            f"  📄 {name}\t{size_str} | {lines_label} {chars_label}{preview_hint}\n"
            f"     读取: file_read(path=\"{name}\")"
        )

    # 图片类
    if file_type == "image":
        w = meta.get("width", 0)
        h = meta.get("height", 0)
        return (
            f"  🖼 {name}\t{size_str} | {w}×{h}px\n"
            f"     读取: file_read(path=\"{name}\")"
        )

    # Parquet
    if file_type == "parquet":
        rows = meta.get("row_count", 0)
        cols = meta.get("col_count", 0)
        col_names = meta.get("columns", [])
        cols_hint = f" | 列: {', '.join(col_names[:5])}" if col_names else ""
        if len(col_names) > 5:
            cols_hint += ", ..."
        return (
            f"  📊 {name}\t{size_str} | {rows:,}行×{cols}列{cols_hint}\n"
            f"     读取: pd.read_parquet('{name}')"
        )

    return f"  {name}\t{size_str}"


def _fmt_size(size: Optional[int]) -> str:
    """格式化文件大小"""
    if not size:
        return "未知大小"
    if size < 1024:
        return f"{size}B"
    if size < 1024 * 1024:
        return f"{size / 1024:.1f}KB"
    return f"{size / 1024 / 1024:.1f}MB"
